import telebot
import google.generativeai as genai
import os
import PyPDF2
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image

# 1. API ve Token Ayarları
TELEGRAM_TOKEN = "SENIN_TELEGRAM_BOT_TOKENIN_BURAYA"
GEMINI_API_KEY = "SENIN_GEMINI_API_ANAHTARIN_BURAYA"

bot = telebot.TeleBot(TELEGRAM_TOKEN)
genai.configure(api_key=GEMINI_API_KEY)

# Hızlı ve vizyon yeteneği olan modeli seçiyoruz
model = genai.GenerativeModel('gemini-1.5-flash')

# Bota /start yazıldığında verilecek cevap
@bot.message_handler(commands=['start'])
def send_welcome(message):
    bot.reply_to(message, "Selam Haşmetli Taha! Ben senin kişisel tıp, acil yardım ve KPSS asistanınım. Sınav notu çıkarmam için bana PDF veya PPTX (slayt) gönder.")

# PDF veya PPTX geldiğinde çalışacak ana kod
@bot.message_handler(content_types=['document'])
def handle_docs(message):
    try:
        file_name = message.document.file_name
        file_info = bot.get_file(message.document.file_id)
        downloaded_file = bot.download_file(file_info.file_path)
        
        # Yapay Zekaya Verilen Kalıcı Uzmanlık Talimatı
        saglik_talimati = "Sen uzman bir paramedik eğitmeni ve KPSS tıbbi bilimler uzmanısın. Gönderilen bu içerikten anatomi, fizyoloji ve acil bakım gibi konularda sınavda çıkabilecek en kritik noktaları, terimleri ve hap bilgileri maddeler halinde çıkar."

        # EĞER DOSYA PPTX (SLAYT) İSE:
        if file_name.endswith('.pptx'):
            bot.reply_to(message, "Slayt geldi kanka! Hemen okuyup paramedik notlarını çıkarıyorum...")
            gecici_dosya = "gecici_sunum.pptx"
            with open(gecici_dosya, 'wb') as new_file:
                new_file.write(downloaded_file)
            
            # Slayt içindeki yazıları çekiyoruz
            prs = Presentation(gecici_dosya)
            metin = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        metin += shape.text + "\n"
                        
            # Yapay zekaya özetletiyoruz
            response = model.generate_content(f"{saglik_talimati}\n\nİçerik: {metin[:15000]}")
            bot.reply_to(message, response.text)
            os.remove(gecici_dosya)

        # EĞER DOSYA PDF İSE:
        elif file_name.endswith('.pdf'):
            bot.reply_to(message, "PDF alındı kanka! Belgeyi inceliyorum...")
            gecici_dosya = "gecici_belge.pdf
            
