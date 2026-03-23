import telebot
import google.generativeai as genai
import os
import PyPDF2
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image

TELEGRAM_TOKEN = os.environ.get("TELEGRAM_TOKEN")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")

bot = telebot.TeleBot(TELEGRAM_TOKEN)
genai.configure(api_key=GEMINI_API_KEY)

# OTOMATİK MODEL SEÇİCİ (Google'ın sunucusundan çalışan güncel modeli kendi bulur)
calisan_model = None
try:
    for m in genai.list_models():
        if 'generateContent' in m.supported_generation_methods:
            calisan_model = m.name
            if 'flash' in calisan_model.lower() or 'pro' in calisan_model.lower():
                break
except Exception as e:
    calisan_model = 'gemini-1.5-flash' # Ne olursa olsun yedekte dursun

if calisan_model is None:
    calisan_model = 'gemini-1.5-flash'

model = genai.GenerativeModel(calisan_model)

@bot.message_handler(commands=['start'])
def send_welcome(message):
    bot.reply_to(message, f"Selam Haşmetli Taha! Sınav notu çıkarmam için bana PDF veya PPTX (slayt) gönder. (Botun bağlandığı yapay zeka: {calisan_model})")

@bot.message_handler(content_types=['document'])
def handle_docs(message):
    try:
        file_name = message.document.file_name
        file_info = bot.get_file(message.document.file_id)
        downloaded_file = bot.download_file(file_info.file_path)
        
        saglik_talimati = "Sen uzman bir paramedik eğitmeni ve KPSS tıbbi bilimler uzmanısın. Gönderilen bu içerikten anatomi, fizyoloji ve acil bakım gibi konularda sınavda çıkabilecek en kritik noktaları, terimleri ve hap bilgileri maddeler halinde çıkar. Karşındaki öğrencinin hedefi atanmak için 95 üstü puan almak, özetlerini bu ciddiyetle hazırla!"

        if file_name.endswith('.pptx'):
            bot.reply_to(message, "Slayt geldi kanka! Hemen okuyup paramedik notlarını çıkarıyorum...")
            gecici_dosya = "gecici_sunum.pptx"
            with open(gecici_dosya, 'wb') as new_file:
                new_file.write(downloaded_file)
            
            prs = Presentation(gecici_dosya)
            metin = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        metin += shape.text + "\n"
                        
            response = model.generate_content(f"{saglik_talimati}\n\nİçerik: {metin[:15000]}")
            bot.reply_to(message, response.text)
            os.remove(gecici_dosya)

        elif file_name.endswith('.pdf'):
            bot.reply_to(message, "PDF alındı kanka! Belgeyi inceliyorum...")
            gecici_dosya = "gecici_belge.pdf"
            with open(gecici_dosya, 'wb') as new_file:
                new_file.write(downloaded_file)
            
            metin = ""
            with open(gecici_dosya, "rb") as pdf_file:
                reader = PyPDF2.PdfReader(pdf_file)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        metin += extracted + "\n"
            
            if len(metin.strip()) < 50:
                bot.reply_to(message, "Bu PDF resim formatında kanka, görüntü işleme (OCR) ile okuyup not çıkarıyorum...")
                pages = convert_from_path(gecici_dosya, 200)
                resim_yolu = 'gecici_sayfa_1.jpg'
                pages[0].save(resim_yolu, 'JPEG')
                
                img = Image.open(resim_yolu)
                response = model.generate_content([saglik_talimati, img])
                bot.reply_to(message, response.text)
                os.remove(resim_yolu)
            
            else:
                response = model.generate_content(f"{saglik_talimati}\n\nİçerik: {metin[:15000]}")
                bot.reply_to(message, response.text)
            
            os.remove(gecici_dosya)

        else:
            bot.reply_to(message, "Kanka sadece .pdf veya .pptx formatında dosya gönderebilirsin!")

    except Exception as e:
        bot.reply_to(message, f"Bir hata oluştu kanka, şuna bir bak: {e}")

bot.polling(none_stop=True)
